from __future__ import annotations

import os
import re
from pathlib import Path
from typing import Any

from .errors import DSLRuntimeError


def _compile_regex(pattern: str, ignore_case: bool = False) -> re.Pattern[str]:
    if not isinstance(pattern, str):
        raise DSLRuntimeError("Regex pattern must be a string")
    flags = re.IGNORECASE if ignore_case else 0
    try:
        return re.compile(pattern, flags)
    except re.error as exc:
        raise DSLRuntimeError(f"Invalid regex pattern: {exc}") from exc


def _render_relative_path(path: Path, cwd: Path) -> str:
    try:
        return Path(os.path.relpath(path.resolve(), cwd.resolve())).as_posix()
    except ValueError:
        return path.resolve().as_posix()


class DSLFile:
    def __init__(
        self,
        path: Path,
        text_chunk_lines: int = 80,
        display_root: Path | None = None,
    ) -> None:
        self.path = path
        self._text_chunk_lines = text_chunk_lines
        self.display_root = (display_root or Path.cwd()).resolve()
        self._chunks_cache: list[str] | None = None
        self._chunks_loaded_from_db = False

    def __repr__(self) -> str:
        return f"File('{self._display_path()}')"

    def __str__(self) -> str:
        return self._display_path()

    def _display_path(self, path: Path | None = None) -> str:
        target = path if path is not None else self.path
        return _render_relative_path(target, self.display_root)

    def read(self, pages=None):
        chunks = self._chunks()
        if pages is None:
            return "\n\n".join(chunks)
        selected = self._normalize_pages(pages, total_pages=len(chunks))
        return [chunks[idx - 1] for idx in selected]

    def search(self, pattern: str, ignore_case: bool = False) -> list[int]:
        regex = _compile_regex(pattern, ignore_case=ignore_case)
        matches = []
        for page_index, chunk in enumerate(self._chunks(), start=1):
            if regex.search(chunk):
                matches.append(page_index)
        return matches

    def contains(self, pattern: str, ignore_case: bool = False) -> bool:
        return bool(self.search(pattern, ignore_case=ignore_case))

    def head(self):
        chunks = self._chunks()
        if not chunks:
            return ""
        return chunks[0]

    def tail(self):
        chunks = self._chunks()
        if not chunks:
            return ""
        return chunks[-1]

    def table(self, max_items: int = 50) -> str:
        if not isinstance(max_items, int) or max_items < 1:
            raise DSLRuntimeError("max_items must be a positive integer")

        entries: list[tuple[int, str, int | None]] = []
        suffix = self.path.suffix.lower()
        if self.path.exists():
            if suffix == ".pdf":
                entries = self._read_pdf_outline(max_items=max_items)
            elif suffix == ".docx":
                entries = self._read_docx_outline(max_items=max_items)
            elif suffix == ".pptx":
                entries = self._read_pptx_outline(max_items=max_items)
        if not entries:
            entries = self._extract_toc_entries_from_text(max_items=max_items)
        if not entries:
            return f"No table of contents detected for {self._display_path()}"
        return (
            f"=== Table of contents for file {self._display_path()} ===\n"
            f"{self._format_toc_tree(entries)}"
        )

    def semantic_search(self, query: str, top_k: int = 5) -> list[int]:
        if not isinstance(query, str) or query.strip() == "":
            raise DSLRuntimeError("query must be a non-empty string")
        if not isinstance(top_k, int) or top_k < 1:
            raise DSLRuntimeError("top_k must be a positive integer")

        from .semantic import semantic_search_file_pages

        return semantic_search_file_pages(
            file_path=self.path,
            query=query,
            top_k=top_k,
            display_root=self.display_root,
        )

    def snippets(
        self,
        pattern: str,
        max_results: int = 5,
        context_chars: int = 80,
        ignore_case: bool = False,
    ) -> list[str]:
        if not isinstance(max_results, int) or max_results < 1:
            raise DSLRuntimeError("max_results must be a positive integer")
        if not isinstance(context_chars, int) or context_chars < 0:
            raise DSLRuntimeError("context_chars must be a non-negative integer")
        regex = _compile_regex(pattern, ignore_case=ignore_case)

        snippets: list[str] = []
        for page_index, chunk in enumerate(self._chunks(), start=1):
            for match in regex.finditer(chunk):
                start = max(match.start() - context_chars, 0)
                end = min(match.end() + context_chars, len(chunk))
                excerpt = chunk[start:end].replace("\n", " ").strip()
                snippets.append(f"[page {page_index}] {excerpt}")
                if len(snippets) >= max_results:
                    return snippets
        return snippets

    def _normalize_pages(self, pages, total_pages: int) -> list[int]:
        if isinstance(pages, int):
            pages_values = [pages]
        elif isinstance(pages, list):
            pages_values = pages
        else:
            raise DSLRuntimeError("pages must be an integer or a list of integers")

        normalized: list[int] = []
        for value in pages_values:
            if not isinstance(value, int):
                raise DSLRuntimeError("pages list must contain only integers")
            if value < 1 or value > total_pages:
                raise DSLRuntimeError(
                    f"Page {value} is out of range for {self.path.name} (1..{total_pages})"
                )
            if value not in normalized:
                normalized.append(value)
        return normalized

    def _chunks(self) -> list[str]:
        if self._chunks_cache is not None:
            return self._chunks_cache

        from .semantic import get_file_pages_from_database

        db_chunks = get_file_pages_from_database(self.path, display_root=self.display_root)
        if db_chunks is not None:
            self._chunks_loaded_from_db = True
            chunks = db_chunks
        else:
            suffix = self.path.suffix.lower()
            if suffix == ".pdf":
                chunks = self._read_pdf_pages()
            elif suffix == ".docx":
                chunks = self._read_docx_chunks()
            elif suffix == ".pptx":
                chunks = self._read_pptx_chunks()
            else:
                chunks = self._read_text_chunks()
        self._chunks_cache = chunks or [""]
        return self._chunks_cache

    def _read_pdf_pages(self) -> list[str]:
        try:
            import pymupdf
        except ImportError as exc:
            raise DSLRuntimeError(
                "PyMuPDF is required to read PDF files. Install dependency 'pymupdf'."
            ) from exc

        try:
            with pymupdf.open(str(self.path)) as doc:
                pages = [page.get_text("text").strip() for page in doc]
        except Exception as exc:
            raise DSLRuntimeError(f"Failed to read PDF '{self.path.name}': {exc}") from exc

        return pages or [""]

    def _read_pdf_outline(self, max_items: int) -> list[tuple[int, str, int | None]]:
        try:
            import pymupdf
        except ImportError as exc:
            raise DSLRuntimeError(
                "PyMuPDF is required to read PDF files. Install dependency 'pymupdf'."
            ) from exc

        try:
            with pymupdf.open(str(self.path)) as doc:
                raw_toc = doc.get_toc(simple=True)
        except Exception as exc:
            raise DSLRuntimeError(f"Failed to read PDF outline '{self.path.name}': {exc}") from exc

        entries: list[tuple[int, str, int | None]] = []
        for item in raw_toc:
            if not isinstance(item, (list, tuple)) or len(item) < 3:
                continue
            raw_level, raw_title, raw_page = item[0], item[1], item[2]
            if not isinstance(raw_level, int):
                continue
            level = max(1, raw_level)
            title = str(raw_title).strip()
            if not title:
                continue

            page: int | None = None
            if isinstance(raw_page, int):
                page = raw_page if raw_page >= 1 else None
            elif isinstance(raw_page, float):
                page_candidate = int(raw_page)
                page = page_candidate if page_candidate >= 1 else None

            entries.append((level, title, page))
            if len(entries) >= max_items:
                break
        return entries

    def _read_docx_chunks(self) -> list[str]:
        try:
            import docx
        except ImportError as exc:
            raise DSLRuntimeError(
                "python-docx is required to read DOCX files. Install dependency 'python-docx'."
            ) from exc

        try:
            doc = docx.Document(str(self.path))
        except Exception as exc:
            raise DSLRuntimeError(f"Failed to read DOCX '{self.path.name}': {exc}") from exc

        chunks: list[str] = []
        current_lines: list[str] = []

        def flush_current() -> None:
            if not current_lines:
                return
            chunk = "\n".join(current_lines).strip()
            if chunk:
                chunks.append(chunk)
            current_lines.clear()

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style_name = (para.style.name if para.style else "") or ""
            style_lower = style_name.lower()
            if style_lower.startswith("heading"):
                flush_current()
                current_lines.append(text)
                continue
            current_lines.append(text)

        for table in doc.tables:
            rows: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                flush_current()
                chunks.append("\n".join(rows))

        flush_current()
        return chunks or [""]

    def _read_docx_outline(self, max_items: int) -> list[tuple[int, str, int | None]]:
        try:
            import docx
        except ImportError as exc:
            raise DSLRuntimeError(
                "python-docx is required to read DOCX files. Install dependency 'python-docx'."
            ) from exc

        try:
            doc = docx.Document(str(self.path))
        except Exception as exc:
            raise DSLRuntimeError(f"Failed to read DOCX outline '{self.path.name}': {exc}") from exc

        entries: list[tuple[int, str, int | None]] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style_name = (para.style.name if para.style else "") or ""
            match = re.match(r"heading\s+(\d+)", style_name.strip().lower())
            if not match:
                continue
            level = max(1, int(match.group(1)))
            entries.append((level, text, None))
            if len(entries) >= max_items:
                break
        return entries

    def _read_pptx_chunks(self) -> list[str]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise DSLRuntimeError(
                "python-pptx is required to read PPTX files. Install dependency 'python-pptx'."
            ) from exc

        try:
            presentation = Presentation(str(self.path))
        except Exception as exc:
            raise DSLRuntimeError(f"Failed to read PPTX '{self.path.name}': {exc}") from exc

        chunks: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
                    continue
                text = shape.text.strip()
                if not text:
                    continue
                for line in text.splitlines():
                    stripped = line.strip()
                    if stripped:
                        lines.append(stripped)

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    lines.append("[Notes]")
                    for line in notes_text.splitlines():
                        stripped = line.strip()
                        if stripped:
                            lines.append(stripped)

            if not lines:
                lines.append(f"[Slide {index}]")
            chunks.append("\n".join(lines))
        return chunks or [""]

    def _read_pptx_outline(self, max_items: int) -> list[tuple[int, str, int | None]]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise DSLRuntimeError(
                "python-pptx is required to read PPTX files. Install dependency 'python-pptx'."
            ) from exc

        try:
            presentation = Presentation(str(self.path))
        except Exception as exc:
            raise DSLRuntimeError(f"Failed to read PPTX outline '{self.path.name}': {exc}") from exc

        entries: list[tuple[int, str, int | None]] = []
        for index, slide in enumerate(presentation.slides, start=1):
            title = ""
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
            if not title:
                for shape in slide.shapes:
                    if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
                        continue
                    text = shape.text.strip()
                    if text:
                        title = text.splitlines()[0].strip()
                        break
            if not title:
                title = f"Slide {index}"

            entries.append((1, title, index))
            if len(entries) >= max_items:
                break
        return entries

    def _extract_toc_entries_from_text(self, max_items: int) -> list[tuple[int, str, int | None]]:
        numbered_dotted = re.compile(r"^(\d+(?:\.\d+)*)\s+(.+?)\.{2,}\s*(\d+)$")
        numbered_plain = re.compile(r"^(\d+(?:\.\d+)*)\s+(.+?)\s+(\d+)$")
        titled_dotted = re.compile(r"^(.+?)\.{2,}\s*(\d+)$")
        entries: list[tuple[int, str, int | None]] = []
        seen: set[tuple[int, str, int | None]] = set()

        for chunk in self._chunks()[:8]:
            for raw_line in chunk.splitlines():
                line = raw_line.strip()
                if len(line) < 8:
                    continue

                level = 1
                title = ""
                page: int | None = None

                match = numbered_dotted.match(line) or numbered_plain.match(line)
                if match:
                    section = match.group(1).strip()
                    body = match.group(2).strip()
                    title = f"{section} {body}".strip()
                    page = int(match.group(3))
                    level = section.count(".") + 1
                else:
                    title_match = titled_dotted.match(line)
                    if title_match:
                        title = title_match.group(1).strip()
                        page = int(title_match.group(2))

                if not title:
                    continue
                key = (level, title, page)
                if key in seen:
                    continue
                seen.add(key)
                entries.append(key)
                if len(entries) >= max_items:
                    return entries
        return entries

    def _format_toc_tree(self, entries: list[tuple[int, str, int | None]]) -> str:
        lines: list[str] = []
        for level, title, page in entries:
            indent = "  " * max(level - 1, 0)
            page_text = f" (p.{page})" if page is not None else ""
            lines.append(f"{indent}{title}{page_text}")
        return "\n".join(lines)

    def _read_text_chunks(self) -> list[str]:
        try:
            text = self.path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            text = self.path.read_text(encoding="utf-8", errors="replace")

        if text == "":
            return [""]

        lines = text.splitlines()
        if not lines:
            return [text]

        chunks: list[str] = []
        for start in range(0, len(lines), self._text_chunk_lines):
            block = "\n".join(lines[start : start + self._text_chunk_lines]).strip()
            chunks.append(block)
        return chunks or [text]


class DSLDirectory:
    def __init__(
        self,
        path: Path,
        recursive: bool = True,
        display_root: Path | None = None,
    ) -> None:
        self.display_root = (display_root or Path.cwd()).resolve()
        if path.exists():
            if not path.is_dir():
                raise DSLRuntimeError(f"Path is not a directory: {self._display_path(path)}")
        elif not self._has_db_backed_files(path):
            raise DSLRuntimeError(f"Directory does not exist: {self._display_path(path)}")
        self.path = path
        self.recursive = recursive

    def __repr__(self) -> str:
        return f"Directory('{self._display_path()}')"

    def __str__(self) -> str:
        return self._display_path()

    def __len__(self) -> int:
        return len(self._iter_file_paths(self.recursive))

    def _display_path(self, path: Path | None = None) -> str:
        target = path if path is not None else self.path
        return _render_relative_path(target, self.display_root)

    def __iter__(self):
        files = self._iter_file_paths(self.recursive)
        for file_path in files:
            yield DSLFile(file_path, display_root=self.display_root)

    def files(self, recursive: bool | None = None) -> list[DSLFile]:
        if recursive is None:
            recursive = self.recursive
        return [
            DSLFile(path, display_root=self.display_root)
            for path in self._iter_file_paths(recursive)
        ]

    def tree(self, max_depth: int = 5, max_entries: int = 500) -> str:
        if not isinstance(max_depth, int) or max_depth < 0:
            raise DSLRuntimeError("max_depth must be a non-negative integer")
        if not isinstance(max_entries, int) or max_entries < 1:
            raise DSLRuntimeError("max_entries must be a positive integer")

        from .semantic import get_directory_file_paths_from_database

        db_paths = get_directory_file_paths_from_database(
            self.path,
            recursive=True,
            display_root=self.display_root,
        )
        if db_paths is not None:
            return self._render_tree_from_paths(db_paths, max_depth=max_depth, max_entries=max_entries)

        lines: list[str] = [f"{self._display_path()}/"]
        emitted = 1
        truncated = False

        def walk(current: Path, depth: int) -> bool:
            nonlocal emitted, truncated
            if depth >= max_depth:
                return True

            try:
                entries = sorted(
                    list(current.iterdir()),
                    key=lambda p: (not p.is_dir(), p.name.lower()),
                )
            except OSError as exc:
                lines.append(f"{'  ' * (depth + 1)}[unreadable: {exc}]")
                emitted += 1
                return emitted < max_entries

            for entry in entries:
                if emitted >= max_entries:
                    truncated = True
                    return False
                label = f"{entry.name}/" if entry.is_dir() else entry.name
                lines.append(f"{'  ' * (depth + 1)}{label}")
                emitted += 1
                if entry.is_dir():
                    keep_going = walk(entry, depth + 1)
                    if not keep_going:
                        return False
            return True

        walk(self.path, 0)
        if truncated:
            lines.append(f"... truncated after {max_entries} entries")
        return "\n".join(lines)

    def _render_tree_from_paths(
        self,
        paths: list[Path],
        *,
        max_depth: int,
        max_entries: int,
    ) -> str:
        root: dict[str, dict[str, Any] | set[str]] = {
            "dirs": {},
            "files": set(),
        }

        for path in paths:
            try:
                relative_path = path.relative_to(self.path)
            except ValueError:
                continue
            if not relative_path.parts:
                continue

            node = root
            for part in relative_path.parts[:-1]:
                dirs = node["dirs"]
                if not isinstance(dirs, dict):
                    break
                node = dirs.setdefault(part, {"dirs": {}, "files": set()})
            else:
                files = node["files"]
                if isinstance(files, set):
                    files.add(relative_path.parts[-1])

        lines: list[str] = [f"{self._display_path()}/"]
        emitted = 1
        truncated = False

        def walk(node: dict[str, dict[str, Any] | set[str]], depth: int) -> bool:
            nonlocal emitted, truncated
            if depth >= max_depth:
                return True

            dirs = node["dirs"]
            files = node["files"]
            if not isinstance(dirs, dict) or not isinstance(files, set):
                return True

            dir_items = sorted(dirs.items(), key=lambda item: item[0].lower())
            file_items = sorted(files, key=lambda item: item.lower())

            for name, child in dir_items:
                if emitted >= max_entries:
                    truncated = True
                    return False
                lines.append(f"{'  ' * (depth + 1)}{name}/")
                emitted += 1
                if isinstance(child, dict):
                    if not walk(child, depth + 1):
                        return False

            for name in file_items:
                if emitted >= max_entries:
                    truncated = True
                    return False
                lines.append(f"{'  ' * (depth + 1)}{name}")
                emitted += 1
            return True

        walk(root, 0)
        if truncated:
            lines.append(f"... truncated after {max_entries} entries")
        return "\n".join(lines)

    def search(
        self,
        pattern: str,
        scope: str = "name",
        in_content: bool = False,
        recursive: bool | None = None,
        ignore_case: bool = False,
    ) -> list[DSLFile]:
        if recursive is None:
            recursive = self.recursive
        if in_content:
            scope = "content"
        if scope not in {"name", "content", "both"}:
            raise DSLRuntimeError("scope must be one of: 'name', 'content', 'both'")

        regex = _compile_regex(pattern, ignore_case=ignore_case)
        files = [
            DSLFile(path, display_root=self.display_root)
            for path in self._iter_file_paths(recursive)
        ]
        matches: list[DSLFile] = []
        for file in files:
            relative = file.path.relative_to(self.path).as_posix()
            name_match = bool(regex.search(file.path.name) or regex.search(relative))
            content_match = False
            if scope in {"content", "both"}:
                content_match = file.contains(pattern, ignore_case=ignore_case)

            if scope == "name" and name_match:
                matches.append(file)
            elif scope == "content" and content_match:
                matches.append(file)
            elif scope == "both" and (name_match or content_match):
                matches.append(file)
        return matches

    def _iter_file_paths(self, recursive: bool) -> list[Path]:
        from .semantic import get_directory_file_paths_from_database

        db_paths = get_directory_file_paths_from_database(
            self.path,
            recursive=recursive,
            display_root=self.display_root,
        )
        if db_paths is not None:
            return db_paths

        if recursive:
            paths = [path for path in self.path.rglob("*") if path.is_file()]
        else:
            paths = [path for path in self.path.glob("*") if path.is_file()]
        paths.sort(key=lambda p: p.as_posix())
        return paths

    def _has_db_backed_files(self, path: Path) -> bool:
        from .semantic import get_directory_file_paths_from_database

        db_paths = get_directory_file_paths_from_database(
            path,
            recursive=True,
            display_root=self.display_root,
        )
        return bool(db_paths)

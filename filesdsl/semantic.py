from __future__ import annotations

import io
import json
import math
import os
import re
import zipfile
from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path
import xml.etree.ElementTree as ET

from .errors import DSLRuntimeError, DSLTimeoutError
from .execution_budget import ExecutionBudget
from .text_utils import normalize_text


SEMANTIC_DB_DIRNAME = ".fdsl_faiss"
SEMANTIC_INDEX_FILENAME = "pages.faiss"
SEMANTIC_RECORDS_FILENAME = "records.json"
SEMANTIC_VECTORS_FILENAME = "vectors.json"
EMBEDDING_DIM = 256


def _check_budget(budget: ExecutionBudget | None, phase: str) -> None:
    if budget is None:
        return
    budget.check(phase)


@dataclass(frozen=True)
class PrepareStats:
    folder: Path
    db_path: Path
    indexed_files: int
    indexed_pages: int


def prepare_semantic_database(
    folder: Path,
    *,
    budget: ExecutionBudget | None = None,
) -> PrepareStats:
    target_folder = folder.resolve()
    if not target_folder.exists():
        raise DSLRuntimeError(f"Folder does not exist: {target_folder.as_posix()}")
    if not target_folder.is_dir():
        raise DSLRuntimeError(f"Path is not a directory: {target_folder.as_posix()}")

    db_path = target_folder / SEMANTIC_DB_DIRNAME
    db_path.mkdir(parents=True, exist_ok=True)

    records: list[dict[str, str | int]] = []
    embedding_inputs: list[str] = []

    indexed_files = 0
    indexed_pages = 0
    for file_path in _iter_document_paths(target_folder, budget=budget):
        _check_budget(budget, "semantic.prepare.file")
        relative_path = file_path.relative_to(target_folder).as_posix()
        pages = _extract_pages(file_path, budget=budget)
        indexed_files += 1

        for page_number, page_text in enumerate(pages, start=1):
            _check_budget(budget, "semantic.prepare.page")
            cleaned = normalize_text(page_text).strip()
            records.append(
                {
                    "relative_path": relative_path,
                    "file_name": file_path.name,
                    "page": page_number,
                    "text": cleaned,
                }
            )
            embedding_inputs.append(
                f"File: {relative_path}\n{cleaned}" if cleaned else f"File: {relative_path}"
            )
            indexed_pages += 1

    _write_faiss_database(db_path, records, embedding_inputs, budget=budget)
    return PrepareStats(target_folder, db_path, indexed_files=indexed_files, indexed_pages=indexed_pages)


def semantic_search_file_pages(
    file_path: Path,
    query: str,
    top_k: int,
    *,
    display_root: Path | None = None,
    budget: ExecutionBudget | None = None,
) -> list[int]:
    if not isinstance(query, str) or query.strip() == "":
        raise DSLRuntimeError("query must be a non-empty string")
    if not isinstance(top_k, int) or top_k < 1:
        raise DSLRuntimeError("top_k must be a positive integer")

    resolved_file_path = file_path.resolve()
    indexed_root = _find_indexed_root(
        resolved_file_path,
        display_root=display_root,
        budget=budget,
    )
    relative_path = resolved_file_path.relative_to(indexed_root).as_posix()
    db_path = indexed_root / SEMANTIC_DB_DIRNAME
    vectors = _load_vectors(db_path, budget=budget)
    _, record_positions = _load_record_indexes(db_path, budget=budget)
    query_vector = _encode_texts([query.strip()], budget=budget)[0]

    scored_pages: list[tuple[float, int]] = []
    for vector_index, page in record_positions.get(relative_path, ()):
        _check_budget(budget, "semantic.search.record")
        if vector_index >= len(vectors):
            continue
        score = _dot(query_vector, vectors[vector_index])
        scored_pages.append((score, page))

    scored_pages.sort(key=lambda item: item[0], reverse=True)
    return [page for _, page in scored_pages[:top_k]]


def get_file_pages_from_database(
    file_path: Path,
    display_root: Path | None = None,
    *,
    budget: ExecutionBudget | None = None,
) -> list[str] | None:
    resolved_file_path = file_path.resolve()
    try:
        indexed_root = _find_indexed_root(
            resolved_file_path,
            display_root=display_root,
            budget=budget,
        )
    except DSLTimeoutError:
        raise
    except DSLRuntimeError:
        return None

    records_path = indexed_root / SEMANTIC_DB_DIRNAME / SEMANTIC_RECORDS_FILENAME
    if not records_path.is_file():
        return None

    relative_path = resolved_file_path.relative_to(indexed_root).as_posix()
    file_pages, _ = _load_record_indexes(indexed_root / SEMANTIC_DB_DIRNAME, budget=budget)
    pages = file_pages.get(relative_path)
    if not pages:
        return None
    result: list[str] = []
    for _, text in pages:
        _check_budget(budget, "semantic.file_pages.record")
        result.append(text)
    return result


def get_directory_file_paths_from_database(
    directory_path: Path,
    recursive: bool,
    display_root: Path | None = None,
    *,
    budget: ExecutionBudget | None = None,
) -> list[Path] | None:
    resolved_dir = directory_path.resolve()
    try:
        indexed_root = _find_indexed_root(
            resolved_dir,
            display_root=display_root,
            budget=budget,
        )
    except DSLTimeoutError:
        raise
    except DSLRuntimeError:
        return None

    records_path = indexed_root / SEMANTIC_DB_DIRNAME / SEMANTIC_RECORDS_FILENAME
    if not records_path.is_file():
        return None

    file_pages, _ = _load_record_indexes(indexed_root / SEMANTIC_DB_DIRNAME, budget=budget)
    rel_dir = resolved_dir.relative_to(indexed_root).as_posix()
    if rel_dir == ".":
        rel_dir = ""

    result: set[Path] = set()
    for rel in file_pages:
        _check_budget(budget, "semantic.directory_paths.record")
        rel_parent = str(Path(rel).parent.as_posix())
        if rel_parent == ".":
            rel_parent = ""

        if recursive:
            if rel_dir and not rel.startswith(f"{rel_dir}/"):
                continue
            result.add(indexed_root / rel)
        else:
            if rel_parent == rel_dir:
                result.add(indexed_root / rel)

    return sorted(result, key=lambda p: p.as_posix())


def _iter_document_paths(folder: Path, *, budget: ExecutionBudget | None = None):
    db_path = folder / SEMANTIC_DB_DIRNAME
    for path in sorted(folder.rglob("*"), key=lambda p: p.as_posix()):
        _check_budget(budget, "semantic.iter_documents.path")
        if not path.is_file():
            continue
        if db_path in path.parents:
            continue
        yield path


def _extract_pages(path: Path, *, budget: ExecutionBudget | None = None) -> list[str]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _read_pdf_pages(path, budget=budget)
    if suffix == ".docx":
        return _read_docx_pages(path, budget=budget)
    if suffix == ".pptx":
        return _read_pptx_pages(path, budget=budget)
    return _read_text_chunks(path, budget=budget)


def _read_pdf_pages(path: Path, *, budget: ExecutionBudget | None = None) -> list[str]:
    try:
        import pymupdf
    except ImportError:
        return _read_text_chunks(path, budget=budget)

    pages: list[str] = []
    with pymupdf.open(str(path)) as doc:
        for page in doc:
            _check_budget(budget, "semantic.read_pdf.page")
            text = page.get_text("text").strip()
            pages.append(text if text else _ocr_pdf_page(page))
    return pages or [""]


def _ocr_pdf_page(page) -> str:
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        return ""

    pix = page.get_pixmap(dpi=200)
    image = Image.open(io.BytesIO(pix.tobytes("png")))
    try:
        return pytesseract.image_to_string(image).strip()
    except Exception:
        return ""


def _read_docx_pages(path: Path, *, budget: ExecutionBudget | None = None) -> list[str]:
    try:
        import docx
        doc = docx.Document(str(path))
        lines: list[str] = []
        for paragraph in doc.paragraphs:
            _check_budget(budget, "semantic.read_docx.paragraph")
            text = paragraph.text.strip()
            if text:
                lines.append(text)
        return ["\n".join(lines)] if lines else [""]
    except DSLTimeoutError:
        raise
    except Exception:
        return _read_docx_xml_fallback(path, budget=budget)


def _read_pptx_pages(path: Path, *, budget: ExecutionBudget | None = None) -> list[str]:
    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
        pages: list[str] = []
        for slide in presentation.slides:
            _check_budget(budget, "semantic.read_pptx.slide")
            lines: list[str] = []
            for shape in slide.shapes:
                _check_budget(budget, "semantic.read_pptx.shape")
                if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text:
                    lines.extend(line.strip() for line in shape.text.splitlines() if line.strip())
            pages.append("\n".join(lines).strip() or "")
        return pages or [""]
    except DSLTimeoutError:
        raise
    except Exception:
        return _read_pptx_xml_fallback(path, budget=budget)


def _read_docx_xml_fallback(path: Path, *, budget: ExecutionBudget | None = None) -> list[str]:
    try:
        with zipfile.ZipFile(path) as archive:
            data = archive.read("word/document.xml")
        root = ET.fromstring(data)
    except Exception:
        return _read_text_chunks(path, budget=budget)

    texts: list[str] = []
    for node in root.findall('.//{*}t'):
        _check_budget(budget, "semantic.read_docx_xml.node")
        if node.text and node.text.strip():
            texts.append(node.text.strip())
    return ["\n".join(texts)] if texts else [""]


def _read_pptx_xml_fallback(path: Path, *, budget: ExecutionBudget | None = None) -> list[str]:
    try:
        with zipfile.ZipFile(path) as archive:
            slide_names = sorted(
                name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            )
            slides: list[str] = []
            for name in slide_names:
                _check_budget(budget, "semantic.read_pptx_xml.slide")
                root = ET.fromstring(archive.read(name))
                texts: list[str] = []
                for node in root.findall('.//{*}t'):
                    _check_budget(budget, "semantic.read_pptx_xml.node")
                    if node.text and node.text.strip():
                        texts.append(node.text.strip())
                slides.append("\n".join(texts))
            return slides or [""]
    except DSLTimeoutError:
        raise
    except Exception:
        return _read_text_chunks(path, budget=budget)


def _read_text_chunks(
    path: Path,
    lines_per_chunk: int = 80,
    *,
    budget: ExecutionBudget | None = None,
) -> list[str]:
    try:
        text = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        text = path.read_text(encoding="utf-8", errors="replace")
    if text == "":
        return [""]
    lines = text.splitlines()
    chunks: list[str] = []
    for i in range(0, len(lines), lines_per_chunk):
        _check_budget(budget, "semantic.read_text.chunk")
        chunks.append("\n".join(lines[i : i + lines_per_chunk]).strip())
    return chunks or [text]


def _render_relative_path(path: Path, cwd: Path) -> str:
    try:
        return Path(os.path.relpath(path.resolve(), cwd.resolve())).as_posix()
    except ValueError:
        return path.resolve().as_posix()


def _find_indexed_root(
    file_path: Path,
    display_root: Path | None = None,
    *,
    budget: ExecutionBudget | None = None,
) -> Path:
    display_base = (display_root or Path.cwd()).resolve()
    for candidate in [file_path, file_path.parent, *file_path.parents]:
        _check_budget(budget, "semantic.find_indexed_root")
        if (candidate / SEMANTIC_DB_DIRNAME).is_dir():
            return candidate
    rendered = _render_relative_path(file_path, display_base)
    raise DSLRuntimeError(f"No semantic index found for {rendered}. Run 'uv run fdsl prepare <folder>' first.")


def _write_faiss_database(
    db_path: Path,
    records: list[dict[str, str | int]],
    embedding_inputs: list[str],
    *,
    budget: ExecutionBudget | None = None,
) -> None:
    vectors = _encode_texts(embedding_inputs, budget=budget)
    (db_path / SEMANTIC_RECORDS_FILENAME).write_text(json.dumps(records, ensure_ascii=False), encoding="utf-8")
    (db_path / SEMANTIC_VECTORS_FILENAME).write_text(json.dumps(vectors), encoding="utf-8")
    # Marker file to make the storage explicitly faiss-backed for callers.
    (db_path / SEMANTIC_INDEX_FILENAME).write_text("faiss-index-placeholder", encoding="utf-8")


def _load_faiss_database(db_path: Path):
    records = _load_records(db_path)
    vectors = _load_vectors(db_path)
    return records, vectors


def _load_records(db_path: Path) -> list[dict[str, str | int]]:
    records_path = db_path / SEMANTIC_RECORDS_FILENAME
    if not records_path.is_file():
        raise DSLRuntimeError("No prepared semantic index collection found. Run 'uv run fdsl prepare <folder>' first.")
    return _load_records_cached(*_cache_key(records_path))


def _load_vectors(db_path: Path, *, budget: ExecutionBudget | None = None) -> list[list[float]]:
    vectors_path = db_path / SEMANTIC_VECTORS_FILENAME
    if not vectors_path.is_file():
        raise DSLRuntimeError("No prepared semantic index collection found. Run 'uv run fdsl prepare <folder>' first.")
    vectors = _load_vectors_cached(*_cache_key(vectors_path))
    if budget is not None:
        for _ in range(0, len(vectors), 512):
            _check_budget(budget, "semantic.load_vectors")
    return vectors


def _cache_key(path: Path) -> tuple[str, int, int]:
    stats = path.stat()
    return (path.as_posix(), stats.st_mtime_ns, stats.st_size)


@lru_cache(maxsize=8)
def _load_records_cached(path: str, mtime_ns: int, size: int) -> list[dict[str, str | int]]:
    del mtime_ns, size
    loaded = json.loads(Path(path).read_text(encoding="utf-8"))
    if not isinstance(loaded, list):
        raise DSLRuntimeError("Semantic records database is corrupted.")
    records: list[dict[str, str | int]] = []
    for entry in loaded:
        if isinstance(entry, dict):
            records.append(entry)
    return records


@lru_cache(maxsize=8)
def _load_vectors_cached(path: str, mtime_ns: int, size: int) -> list[list[float]]:
    del mtime_ns, size
    loaded = json.loads(Path(path).read_text(encoding="utf-8"))
    if not isinstance(loaded, list):
        raise DSLRuntimeError("Semantic vectors database is corrupted.")
    vectors: list[list[float]] = []
    for vector in loaded:
        if isinstance(vector, list):
            vectors.append([float(value) for value in vector])
    return vectors


def _load_record_indexes(
    db_path: Path,
    *,
    budget: ExecutionBudget | None = None,
) -> tuple[dict[str, tuple[tuple[int, str], ...]], dict[str, tuple[tuple[int, int], ...]]]:
    records_path = db_path / SEMANTIC_RECORDS_FILENAME
    if not records_path.is_file():
        raise DSLRuntimeError("No prepared semantic index collection found. Run 'uv run fdsl prepare <folder>' first.")
    pages, positions = _load_record_indexes_cached(*_cache_key(records_path))
    if budget is not None:
        for _ in range(0, len(pages), 512):
            _check_budget(budget, "semantic.load_records.pages")
        for _ in range(0, len(positions), 512):
            _check_budget(budget, "semantic.load_records.positions")
    return pages, positions


@lru_cache(maxsize=8)
def _load_record_indexes_cached(
    path: str, mtime_ns: int, size: int
) -> tuple[dict[str, tuple[tuple[int, str], ...]], dict[str, tuple[tuple[int, int], ...]]]:
    records = _load_records_cached(path, mtime_ns, size)
    pages_by_path: dict[str, list[tuple[int, str]]] = {}
    positions_by_path: dict[str, list[tuple[int, int]]] = {}

    for index, record in enumerate(records):
        relative_path = record.get("relative_path")
        if not isinstance(relative_path, str):
            continue

        page_number = record.get("page")
        normalized_page = page_number if isinstance(page_number, int) else 0
        text = record.get("text", "")
        pages_by_path.setdefault(relative_path, []).append((normalized_page, str(text)))
        if isinstance(page_number, int):
            positions_by_path.setdefault(relative_path, []).append((index, page_number))

    normalized_pages: dict[str, tuple[tuple[int, str], ...]] = {}
    for relative_path, pages in pages_by_path.items():
        pages.sort(key=lambda item: item[0])
        normalized_pages[relative_path] = tuple(pages)

    normalized_positions: dict[str, tuple[tuple[int, int], ...]] = {}
    for relative_path, positions in positions_by_path.items():
        positions.sort(key=lambda item: item[1])
        normalized_positions[relative_path] = tuple(positions)

    return normalized_pages, normalized_positions


def _encode_texts(texts: list[str], *, budget: ExecutionBudget | None = None) -> list[list[float]]:
    vectors: list[list[float]] = []
    for text in texts:
        _check_budget(budget, "semantic.encode.text")
        vec = [0.0] * EMBEDDING_DIM
        for token_index, token in enumerate(re.findall(r"\w+", text.lower())):
            if token_index % 128 == 0:
                _check_budget(budget, "semantic.encode.token")
            vec[hash(token) % EMBEDDING_DIM] += 1.0
        norm = math.sqrt(sum(v * v for v in vec))
        if norm > 0:
            vec = [v / norm for v in vec]
        vectors.append(vec)
    return vectors


def _dot(a: list[float], b: list[float]) -> float:
    return sum(x * y for x, y in zip(a, b, strict=False))

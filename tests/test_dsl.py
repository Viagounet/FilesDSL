from __future__ import annotations

from concurrent.futures import ThreadPoolExecutor
from io import StringIO
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from filesdsl.errors import DSLRuntimeError, DSLSyntaxError
from filesdsl.interpreter import execute_fdsl, run_script


class FilesDSLTests(unittest.TestCase):
    def test_execute_fdsl_from_python_string(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            code = "x = 1\npages = [2:4]\nprint(x)\nprint(pages)\n"
            output = execute_fdsl(code, cwd=root, sandbox_root=root)
            self.assertEqual(output, "1\n[2, 3, 4]\n")

    def test_execute_fdsl_output_isolated_per_call_under_threads(self) -> None:
        root = Path.cwd()

        def run_case(value: int) -> str:
            code = f'value = {value}\nprint("start", value)\nprint([1:3])\nprint("end", value)\n'
            return execute_fdsl(code, cwd=root, sandbox_root=root)

        for _ in range(5):
            with ThreadPoolExecutor(max_workers=16) as executor:
                outputs = list(executor.map(run_case, range(64)))

            for value, output in enumerate(outputs):
                expected = f"start {value}\n[1, 2, 3]\nend {value}\n"
                self.assertEqual(output, expected)

    def test_run_script_writes_prints_to_given_stdout(self) -> None:
        root = Path.cwd()
        output = StringIO()
        run_script('print("alpha", 1)\n', cwd=root, sandbox_root=root, stdout=output)
        self.assertEqual(output.getvalue(), "alpha 1\n")

    def test_range_syntax_in_lists(self) -> None:
        script = "pages = [1, 5:8, 15]\n"
        variables = run_script(script, cwd=Path.cwd(), sandbox_root=Path.cwd())
        self.assertEqual(variables["pages"], [1, 5, 6, 7, 8, 15])

    def test_multiline_list_literal(self) -> None:
        script = """
terms = [
    "a",
    "b",
    "c",
    "d",
]
"""
        variables = run_script(script, cwd=Path.cwd(), sandbox_root=Path.cwd())
        self.assertEqual(variables["terms"], ["a", "b", "c", "d"])

    def test_directory_search_and_file_api(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            (root / "a.txt").write_text("alpha\nbeta\ngamma\n", encoding="utf-8")
            (root / "b.txt").write_text("delta\n", encoding="utf-8")
            (root / "nested").mkdir()
            (root / "nested" / "c.txt").write_text("alpha inside nested\n", encoding="utf-8")

            script = """
docs = Directory(".")
matches = docs.search("a\\.txt$", scope="name")
hits = []
for file in matches:
    if file.contains("alpha"):
        pages = file.search("alpha")
        chunks = file.read(pages=[1])
        first = file.head()
        last = file.tail()
        hits = hits + pages
"""
            variables = run_script(script, cwd=root, sandbox_root=root)
            self.assertEqual(variables["hits"], [1])
            self.assertEqual(len(variables["chunks"]), 1)
            self.assertTrue("alpha" in variables["first"])
            self.assertTrue("gamma" in variables["last"])

            recursive_script = """
docs = Directory(".")
nested = docs.search("c\\.txt$", scope="name")
count = len(nested)
"""
            recursive_vars = run_script(recursive_script, cwd=root, sandbox_root=root)
            self.assertEqual(recursive_vars["count"], 1)

    def test_len_directory_returns_file_count(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            (root / "a.txt").write_text("a", encoding="utf-8")
            (root / "b.txt").write_text("b", encoding="utf-8")
            (root / "nested").mkdir()
            (root / "nested" / "c.txt").write_text("c", encoding="utf-8")

            script = """
docs = Directory(".")
count_recursive = len(docs)
count_flat = len(Directory(".", recursive=false))
"""
            variables = run_script(script, cwd=root, sandbox_root=root)
            self.assertIsInstance(variables["count_recursive"], int)
            self.assertIsInstance(variables["count_flat"], int)
            self.assertEqual(variables["count_recursive"], 3)
            self.assertEqual(variables["count_flat"], 2)

    def test_file_builtin_direct_access(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            (root / "single.txt").write_text("alpha\nbeta\ngamma\n", encoding="utf-8")
            script = """
f = File("single.txt")
has_alpha = f.contains("alpha")
alpha_pages = f.search("alpha")
first = f.head()
"""
            variables = run_script(script, cwd=root, sandbox_root=root)
            self.assertTrue(variables["has_alpha"])
            self.assertEqual(variables["alpha_pages"], [1])
            self.assertIn("alpha", variables["first"])

    def test_semantic_search_file_method(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            (root / "single.txt").write_text("alpha\nbeta\ngamma\n", encoding="utf-8")
            script = """
f = File("single.txt")
hits = f.semantic_search("alpha", top_k=2)
"""
            with patch("filesdsl.semantic.semantic_search_file_pages", return_value=[2, 1]) as search_mock:
                variables = run_script(script, cwd=root, sandbox_root=root)

            self.assertEqual(variables["hits"], [2, 1])
            search_mock.assert_called_once()

    def test_semantic_search_rejects_invalid_top_k(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            (root / "single.txt").write_text("alpha\nbeta\ngamma\n", encoding="utf-8")
            script = """
f = File("single.txt")
hits = f.semantic_search("alpha", top_k=0)
"""
            with self.assertRaises(DSLRuntimeError):
                run_script(script, cwd=root, sandbox_root=root)

    def test_printed_paths_are_relative_to_cwd(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            (root / "single.txt").write_text("alpha\n", encoding="utf-8")
            script = """
d = Directory(".")
f = File("single.txt")
print(d)
print(f)
"""
            output = execute_fdsl(script, cwd=root, sandbox_root=root)
            self.assertEqual(output, ".\nsingle.txt\n")

    def test_sandbox_denies_outside_root(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            script = "docs = Directory('/')\n"
            with self.assertRaises(DSLRuntimeError):
                run_script(script, cwd=root, sandbox_root=root)

    def test_file_builtin_denies_outside_root(self) -> None:
        with tempfile.TemporaryDirectory() as root_dir, tempfile.TemporaryDirectory() as outside_dir:
            root = Path(root_dir)
            outside_file = Path(outside_dir) / "external.txt"
            outside_file.write_text("outside", encoding="utf-8")
            script = f"f = File('{outside_file.as_posix()}')\n"
            with self.assertRaises(DSLRuntimeError):
                run_script(script, cwd=root, sandbox_root=root)

    def test_table_returns_tree_string(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            (root / "toc.txt").write_text(
                "1 Introduction ........ 1\n"
                "1.1 Scope ........ 2\n"
                "2 Methods ........ 5\n",
                encoding="utf-8",
            )
            script = """
docs = Directory(".")
files = docs.search("toc\\.txt$", scope="name")
toc = ""
for file in files:
    toc = file.table()
"""
            variables = run_script(script, cwd=root, sandbox_root=root)
            self.assertIsInstance(variables["toc"], str)
            self.assertIn(
                "=== Table of contents for file toc.txt ===",
                variables["toc"],
            )
            self.assertIn("1 Introduction (p.1)", variables["toc"])
            self.assertIn("  1.1 Scope (p.2)", variables["toc"])
            self.assertIn("2 Methods (p.5)", variables["toc"])

    def test_table_returns_no_toc_message(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            (root / "plain.txt").write_text("hello\nworld\n", encoding="utf-8")
            script = """
docs = Directory(".")
files = docs.search("plain\\.txt$", scope="name")
toc = ""
for file in files:
    toc = file.table()
"""
            variables = run_script(script, cwd=root, sandbox_root=root)
            expected = "No table of contents detected for plain.txt"
            self.assertEqual(variables["toc"], expected)

    def test_directory_tree_output(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            (root / "top.txt").write_text("x", encoding="utf-8")
            (root / "nested").mkdir()
            (root / "nested" / "inside.txt").write_text("y", encoding="utf-8")

            script = """
docs = Directory(".")
tree_text = docs.tree(max_depth=3)
"""
            variables = run_script(script, cwd=root, sandbox_root=root)
            tree_text = variables["tree_text"]
            self.assertEqual(tree_text.splitlines()[0], "./")
            self.assertIn("  nested/", tree_text)
            self.assertIn("    inside.txt", tree_text)
            self.assertIn("  top.txt", tree_text)

    def test_directory_tree_truncation(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            for i in range(10):
                (root / f"file_{i}.txt").write_text("x", encoding="utf-8")
            script = """
docs = Directory(".")
tree_text = docs.tree(max_entries=3)
"""
            variables = run_script(script, cwd=root, sandbox_root=root)
            self.assertIn("truncated after 3 entries", variables["tree_text"])

    def test_syntax_error_reports_location(self) -> None:
        script = "for file in Directory('.')\n    print(file)\n"
        with self.assertRaises(DSLSyntaxError) as context:
            run_script(script, cwd=Path.cwd(), sandbox_root=Path.cwd())
        self.assertEqual(context.exception.line, 1)
        self.assertGreaterEqual(context.exception.column, 1)

    def test_docx_file_methods(self) -> None:
        try:
            from docx import Document
        except ModuleNotFoundError:
            self.skipTest("python-docx is not installed")

        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            path = root / "sample.docx"

            doc = Document()
            doc.add_heading("Overview", level=1)
            doc.add_paragraph("alpha signal in overview")
            doc.add_heading("Details", level=2)
            doc.add_paragraph("beta signal in details")
            doc.save(path)

            script = """
f = File("sample.docx")
all_text = f.read()
selected = f.read(pages=[1, 2])
matches = f.search("alpha")
has_beta = f.contains("beta")
first = f.head()
last = f.tail()
toc = f.table()
snips = f.snippets("alpha", max_results=1)
"""
            variables = run_script(script, cwd=root, sandbox_root=root)

            self.assertIn("Overview", variables["all_text"])
            self.assertIsInstance(variables["selected"], list)
            self.assertEqual(variables["matches"], [1])
            self.assertTrue(variables["has_beta"])
            self.assertIn("Overview", variables["first"])
            self.assertIn("Details", variables["last"])
            self.assertIn("Overview", variables["toc"])
            self.assertIn("  Details", variables["toc"])
            self.assertTrue(variables["snips"])
            self.assertIn("[page 1]", variables["snips"][0])

    def test_pptx_file_methods(self) -> None:
        try:
            from pptx import Presentation
        except ModuleNotFoundError:
            self.skipTest("python-pptx is not installed")

        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            path = root / "sample.pptx"

            presentation = Presentation()
            slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide1.shapes.title.text = "Intro"
            slide1.shapes.placeholders[1].text = "alpha launch notes"

            slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide2.shapes.title.text = "Methods"
            slide2.shapes.placeholders[1].text = "beta evaluation plan"
            presentation.save(path)

            script = """
f = File("sample.pptx")
all_text = f.read()
selected = f.read(pages=[2])
matches = f.search("alpha")
has_beta = f.contains("beta")
first = f.head()
last = f.tail()
toc = f.table()
snips = f.snippets("alpha", max_results=1)
"""
            variables = run_script(script, cwd=root, sandbox_root=root)

            self.assertIn("Intro", variables["all_text"])
            self.assertEqual(variables["matches"], [1])
            self.assertTrue(variables["has_beta"])
            self.assertIsInstance(variables["selected"], list)
            self.assertIn("Methods", variables["selected"][0])
            self.assertIn("Intro", variables["first"])
            self.assertIn("Methods", variables["last"])
            self.assertIn("Intro (p.1)", variables["toc"])
            self.assertIn("Methods (p.2)", variables["toc"])
            self.assertTrue(variables["snips"])
            self.assertIn("[page 1]", variables["snips"][0])


if __name__ == "__main__":
    unittest.main()
